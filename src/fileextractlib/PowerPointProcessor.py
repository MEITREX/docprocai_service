import argparse
import json
from typing import Any

from pptx import Presentation


class PowerPointProcessor:

    def extract_text_from_pptx(self, filename: str) -> dict[Any, Any]:
        prs = Presentation(filename)
        text_runs = {}

        for page, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text.strip():
                    text += shape.text

                    text_runs[page + 1] = text

        with open("test.json", "w", encoding="utf8") as file:
            json.dump(text_runs, file, indent=2, ensure_ascii=False)

        return text_runs


parser = argparse.ArgumentParser()
parser.add_argument("--file")
args = parser.parse_args()
file_path = 'path_to_your_pptx_file.pptx'
processor = PowerPointProcessor()
text_content = processor.extract_text_from_pptx(args.file)
print(text_content)
